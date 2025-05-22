import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter

def load_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

def load_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

def load_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def load_pptx(file_path):
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return "\n".join(texts)

def load_file_text(file_path):
    ext = os.path.splitext(file_path)[-1].lower()
    if ext == ".txt":
        return load_txt(file_path)
    elif ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".docx":
        return load_docx(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

def chunk_text(text, chunk_size=500, chunk_overlap=50):
    splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return splitter.split_text(text)